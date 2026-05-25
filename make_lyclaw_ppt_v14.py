#!/usr/bin/env python3
"""LyClaw PPT v14 —— 宽松排版 · 多图表 · 少文字 · 新增未来方向"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── 调色板 ──
PRI = RGBColor(0xCC, 0x78, 0x5C)  # 主色
DK  = RGBColor(0x18, 0x17, 0x15)
AMB = RGBColor(0xE8, 0xA5, 0x5A)
GRN = RGBColor(0x5D, 0xB8, 0x72)
TEA = RGBColor(0x5D, 0xB8, 0xA6)
BRO = RGBColor(0xA9, 0x58, 0x3E)
BLU = RGBColor(0x5B, 0x8D, 0xC8)
PNK = RGBColor(0xD4, 0x6B, 0x8A)
W   = RGBColor(0xFF, 0xFF, 0xFF)
BG  = RGBColor(0xFA, 0xF9, 0xF5)
TXT = RGBColor(0x3D, 0x3D, 0x3A)
MUT = RGBColor(0x6C, 0x6A, 0x64)
LINE= RGBColor(0xE6, 0xDF, 0xD8)
C = [PRI, TEA, AMB, GRN, BRO, BLU, PNK]
CBG = RGBColor(0xF5, 0xF0, 0xE8)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
MG = Inches(0.6)
CW = SW - 2 * MG

def fill(sl, color=BG):
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = color

def rnd(sl, l, t, w, h, color, r=0.06):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    s.adjustments[0] = r; return s

def arr(sl, l, t, w, h, color):
    s = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background(); return s

def rect(sl, l, t, w, h, c, r=None):
    tp = MSO_SHAPE.ROUNDED_RECTANGLE if r else MSO_SHAPE.RECTANGLE
    s = sl.shapes.add_shape(tp, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()
    if r: s.adjustments[0] = r; return s

def tb(sl, l, t, w, h, text, sz=14, c=TXT, b=False, al=PP_ALIGN.LEFT, fname="Microsoft YaHei"):
    tx = sl.shapes.add_textbox(l, t, w, h)
    tx.text_frame.word_wrap = True
    p = tx.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = b; p.font.name = fname; p.alignment = al; return tx

def mul(sl, l, t, w, h, items, sz=13, c=TXT, sp=Pt(2)):
    tx = sl.shapes.add_textbox(l, t, w, h)
    tx.text_frame.word_wrap = True
    for i, item in enumerate(items):
        p = tx.text_frame.paragraphs[0] if i == 0 else tx.text_frame.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = c
        p.font.name = "Microsoft YaHei"; p.space_after = sp
    return tx

def page_title(sl, num, title, sub=""):
    tb(sl, MG, Inches(0.2), Inches(0.45), Inches(0.3), num, 13, PRI, True)
    tb(sl, MG + Inches(0.4), Inches(0.15), Inches(10), Inches(0.35), title, 26, DK, True)
    if sub: tb(sl, MG, Inches(0.65), CW, Inches(0.2), sub, 11, MUT)

def card(sl, l, t, w, h, accent, title, items, ts=13, bs=11):
    rnd(sl, l, t, w, h, W); rect(sl, l, t, w, Inches(0.04), accent)
    tb(sl, l+Inches(0.12), t+Inches(0.08), w-Inches(0.24), Inches(0.22), title, ts, DK, True)
    mul(sl, l+Inches(0.12), t+Inches(0.35), w-Inches(0.24), h-Inches(0.4), items, bs, TXT, Pt(1))

def tag(sl, l, t, w, h, c, text, sz=9):
    rect(sl, l, t, w, h, c, 0.03)
    tb(sl, l+Inches(0.02), t+Inches(0.01), w-Inches(0.04), h-Inches(0.02), text, sz, W, True, PP_ALIGN.CENTER)

# ════════════════════════  1  封面  ════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); fill(sl, W)
rect(sl, 0, 0, SW, Inches(0.05), PRI)
rect(sl, 0, Inches(0.05), SW, Inches(2.1), DK)
tb(sl, MG, Inches(0.45), Inches(11), Inches(0.7), "LyClaw 技术亮点串讲", 34, W, True)
tb(sl, MG, Inches(1.2), Inches(11), Inches(0.35),
   "JDK 动态代理 · Pipeline · ReAct 引擎 · 反思拓扑 · 工具系统", 14, RGBColor(0xE8, 0xE0, 0xD2))

badges = [("@Agent\n接口", C[0]), ("JDK Proxy\n动态代理", C[1]),
          ("Pipeline\n拓扑排序", C[2]), ("ReAct\n状态机", C[3]),
          ("反思\n策略+组合", C[4]), ("工具\n适配器+外观", C[5])]
for i, (lbl, c) in enumerate(badges):
    x = MG + Inches(i * 2.0)
    rnd(sl, x, Inches(2.5), Inches(1.75), Inches(0.65), c)
    tb(sl, x+Inches(0.04), Inches(2.52), Inches(1.67), Inches(0.6), lbl, 12, W, True, al=PP_ALIGN.CENTER)
    if i < len(badges)-1: arr(sl, x+Inches(1.78), Inches(2.68), Inches(0.2), Inches(0.25), LINE)

tb(sl, MG, Inches(3.6), Inches(11), Inches(0.25),
   "从 @Agent 接口到工具执行 —— 五种经典设计模式协同工作的一条完整链路", 13, TXT, True)
tb(sl, MG, Inches(4.1), Inches(5), Inches(0.25), "软件系统分析与设计 · 2026-05-25", 12, MUT)
tb(sl, MG, Inches(4.8), Inches(11), Inches(0.3), "一条链路 · 五种模式 · 零耦合设计", 16, PRI, True)

# 技术栈标签
techs = ["Java 21", "Spring Boot 3.5", "WebFlux", "React", "SSE 流式推送", "Vue 3.5"]
for i, t in enumerate(techs):
    x = MG + Inches(i * 2.0)
    rect(sl, x, Inches(5.4), Inches(1.75), Inches(0.35), C[i % len(C)], 0.04)
    tb(sl, x+Inches(0.04), Inches(5.43), Inches(1.67), Inches(0.28), t, 11, W, True, al=PP_ALIGN.CENTER)

# ════════════════════════  2  提纲  ════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); fill(sl)
page_title(sl, "02", "汇报提纲")

items = [
    ("01", "项目背景", "通用智能体框架的定位与目标", PRI),
    ("02", "JDK 动态代理", "接口声明式 → 自动 Agent 代理", TEA),
    ("03", "Pipeline 管道", "拓扑排序 · 5 阶段 · 3 种运行模式", AMB),
    ("04", "ReAct 推理引擎", "3 状态状态机 · 多轮循环 · 审批", GRN),
    ("05", "反思拓扑 DAG", "6 原语 · 6 预设 · 6 评估器", BRO),
    ("06", "工具系统", "@Tool 注解 · 适配器 · 3 级沙箱", BLU),
    ("07", "未来演进", "记忆架构 · 多 Agent 调度 · RAG", PNK),
]
for i, (num, title, desc, c) in enumerate(items):
    y = Inches(1.1 + i * 0.82)
    rnd(sl, MG, y, CW, Inches(0.68), W)
    rect(sl, MG, y, Inches(0.06), Inches(0.68), c)
    tag(sl, MG+Inches(0.2), y+Inches(0.14), Inches(0.45), Inches(0.35), c, num)
    tb(sl, MG+Inches(0.85), y+Inches(0.06), Inches(3), Inches(0.25), title, 17, DK, True)
    tb(sl, MG+Inches(0.85), y+Inches(0.38), Inches(8), Inches(0.2), desc, 11, MUT)

# ════════════════════════  3  项目背景  ════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); fill(sl)
page_title(sl, "03", "项目背景",
    "从痛点出发，构建通用、响应式的 AI Agent 框架")

bgs = [
    ("痛点", PRI, [
        "▸ 现有框架配置复杂、启动慢",
        "▸ 缺乏反应式流 (Reactive) 支持",
        "▸ Agent 编排不够灵活",
        "▸ 工具集成和安全管理繁琐",
    ]),
    ("目标", TEA, [
        "▸ 轻量级、零配置自动装配",
        "▸ Pipeline 声明式编排 + 拓扑排序",
        "▸ 支持 ReAct / Reflection / Plan 多模式",
        "▸ 3 级沙箱保障工具执行安全",
    ]),
    ("定位", AMB, [
        "▸ Java 21 + Spring Boot 3.5 WebFlux",
        "▸ JDK 动态代理 (非 CGLIB)",
        "▸ SSE 流式推送前端实时展示",
        "▸ 一条链路 · 五种模式 · 零耦合",
    ]),
]
for i, (title, c, items) in enumerate(bgs):
    x = MG + Inches(i * 4.05)
    rnd(sl, x, Inches(1.5), Inches(3.75), Inches(4.5), W)
    rect(sl, x, Inches(1.5), Inches(3.75), Inches(0.06), c)
    tb(sl, x+Inches(0.15), Inches(1.65), Inches(3.45), Inches(0.3), title, 20, c, True)
    rect(sl, x+Inches(0.15), Inches(2.05), Inches(1.5), Inches(0.02), c)
    mul(sl, x+Inches(0.15), Inches(2.25), Inches(3.45), Inches(3.5), items, 13, TXT, Pt(6))

# ════════════════════════  4  JDK 代理  ════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); fill(sl)
page_title(sl, "05", "JDK 动态代理", "接口声明式 → 自动 Agent 代理")

# 流程图 - 5 steps
steps = [("@Agent\n扫描", C[0]), ("AgentInterface\nProcessor", C[1]),
         ("AgentProxy\nFactory", C[2]), ("Proxy.new\nProxyInstance", C[3]),
         ("Handler\n.invoke", C[4])]
for i, (lbl, c) in enumerate(steps):
    x = MG + Inches(i * 2.35)
    rnd(sl, x, Inches(1.05), Inches(2.0), Inches(0.9), c)
    tb(sl, x+Inches(0.06), Inches(1.08), Inches(1.88), Inches(0.75), lbl, 12, W, True, al=PP_ALIGN.CENTER)
    if i < len(steps)-1: arr(sl, x+Inches(2.04), Inches(1.35), Inches(0.28), Inches(0.3), MUT)

# 代码示例 - 更紧凑
code = ('@Agent(model = "deepseek-v4", provider = "openai")\n'
        'public interface ChatAgent {\n'
        '    @SystemMessage("你好")\n'
        '    Flux<ServerSentEvent<String>>\n'
        '        chat(@UserMessage @V("query") String query);\n'
        '}')
rnd(sl, MG, Inches(2.2), Inches(5.7), Inches(1.5), W)
rect(sl, MG, Inches(2.2), Inches(5.7), Inches(0.04), PRI)
tb(sl, MG+Inches(0.1), Inches(2.3), Inches(2), Inches(0.2), "@Agent 接口", 13, DK, True)
rnd(sl, MG+Inches(0.1), Inches(2.55), Inches(5.5), Inches(1.05), BG)
tb(sl, MG+Inches(0.18), Inches(2.6), Inches(5.35), Inches(0.95), code, 10, TXT, fname="Consolas")

card(sl, MG, Inches(4.0), Inches(5.7), Inches(1.6), TEA, "执行路径选择", [
    "hasExecutionStage ?",
    "  Flux.concat(stages)",
    "  : ReActEngine.executeStream()",
], 12, 11)

rnd(sl, Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.8), W)
rect(sl, Inches(7.2), Inches(2.2), Inches(5.5), Inches(0.04), AMB)
tb(sl, Inches(7.35), Inches(2.3), Inches(5), Inches(0.2), "代理三要素", 14, DK, True)
mul(sl, Inches(7.35), Inches(2.65), Inches(5.2), Inches(4.0), [
    "▸ AgentInterfaceProcessor (BFPP)",
    "  扫描 @Agent → 注册 FactoryBean",
    "",
    "▸ AgentProxyFactory.create()",
    "  @Agent(model/provider/systemPrompt)",
    "  → AgentInvocationHandler",
    "  → Proxy.newProxyInstance(cL, [interface], handler)",
    "  纯 JDK Proxy！不是 CGLIB / Spring AOP",
    "",
    "▸ AgentInvocationHandler.invoke()",
    "  拦截方法 → 解析 @SystemMessage/@UserMessage",
    "  → ChatRequest + AgentContext",
    "  → mode→profile→管线/直连",
], 11, TXT, Pt(2))

tb(sl, MG, Inches(7.1), CW, Inches(0.2), "设计模式：代理 (Proxy)", 10, MUT)

# ════════════════════════  4  Pipeline 架构  ════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); fill(sl)
page_title(sl, "05", "Pipeline 管道架构",
    "ReactivePipelineStage · 拓扑排序 · 5 阶段")

# 5 stages as a connected horizontal flow
stages = [
    ("ContextBuild", C[0]), ("SecurityCheck", C[1]),
    ("PlanExecution", C[2]), ("Reflection\nTopology", C[3]),
    ("Metrics", C[4]),
]
for i, (lbl, c) in enumerate(stages):
    x = MG + Inches(i * 2.45)
    rnd(sl, x, Inches(1.1), Inches(2.2), Inches(0.65), c)
    tb(sl, x+Inches(0.04), Inches(1.15), Inches(2.12), Inches(0.55), lbl, 11, W, True, al=PP_ALIGN.CENTER)
    if i < len(stages)-1: arr(sl, x+Inches(2.24), Inches(1.25), Inches(0.2), Inches(0.3), MUT)

# 一张卡片说明核心
rnd(sl, MG, Inches(2.1), Inches(6.0), Inches(1.5), W)
rect(sl, MG, Inches(2.1), Inches(6.0), Inches(0.04), PRI)
tb(sl, MG+Inches(0.1), Inches(2.18), Inches(2), Inches(0.2), "核心接口与注解", 13, DK, True)
mul(sl, MG+Inches(0.1), Inches(2.45), Inches(5.8), Inches(1.0), [
    "▸ ReactivePipelineStage: execute(ctx), getOrder(), supportsProfile()",
    "▸ @PipelineStage(name, after={X}, before={Y}, group)",
    "▸ PipelineStageBase: sseEvent() / escapeJson() / logJson()",
    "▸ TopologySort (Kahn) + PipelineStageProcessor (BPP) → 自动发现编排",
], 11, TXT, Pt(1))

# 执行路径选择 - 简化
rnd(sl, MG, Inches(3.9), CW, Inches(2.8), W)
rect(sl, MG, Inches(3.9), CW, Inches(0.04), PRI)
tb(sl, MG+Inches(0.12), Inches(3.98), Inches(3), Inches(0.22), "执行路径选择", 14, DK, True)

paths = [("extras.mode\n→ Profile", C[2]), ("supportsProfile\n筛选阶段", C[3]),
         ("有执行阶段?\n→ Flux.concat", C[0]), ("无执行阶段?\n→ ReActEngine", C[4])]
for i, (lbl, c) in enumerate(paths):
    x = MG + Inches(0.3 + i * 3.05)
    rnd(sl, x, Inches(4.5), Inches(2.7), Inches(0.8), c)
    tb(sl, x+Inches(0.06), Inches(4.53), Inches(2.58), Inches(0.7), lbl, 11, W, True, al=PP_ALIGN.CENTER)
    if i < len(paths)-1: arr(sl, x+Inches(2.74), Inches(4.75), Inches(0.3), Inches(0.25), MUT)

tb(sl, MG, Inches(7.2), CW, Inches(0.2), "设计模式：责任链 (CoR) — 声明式编排，按需激活", 10, MUT)

# ════════════════════════  5  3 种运行模式  ════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); fill(sl)
page_title(sl, "05", "Pipeline 运行模式",
    "react · reflection · plan-execute — 不同阶段激活组合")

profs = [
    ("react", PRI, "纯 ReAct 工具调用", "默认模式", [
        "无执行型阶段激活",
        "→ 直接 ReActEngine",
        "适用：简单对话",
    ]),
    ("reflection", TEA, "自评估 DAG 闭环", "extras.mode=reflection", [
        "ContextBuild → SecurityCheck →",
        "ReflectionTopologyStage → Metrics",
        "适用：需要自我纠错",
    ]),
    ("plan-execute", AMB, "任务分解逐步执行", "TaskPlanner", [
        "PlanExecutionStage(order=2)",
        "→ Metrics",
        "适用：复杂多步推理",
    ]),
]
for i, (name, c, subtitle, mode, items) in enumerate(profs):
    x = MG + Inches(i * 4.1)
    rnd(sl, x, Inches(1.1), Inches(3.8), Inches(4.8), W)
    rect(sl, x, Inches(1.1), Inches(3.8), Inches(0.06), c)
    tb(sl, x+Inches(0.15), Inches(1.25), Inches(3.5), Inches(0.25), name, 18, c, True)
    tb(sl, x+Inches(0.15), Inches(1.55), Inches(3.5), Inches(0.2), subtitle, 11, TXT)
    tag(sl, x+Inches(0.15), Inches(1.85), Inches(1.5), Inches(0.25), c, mode)
    rect(sl, x+Inches(0.15), Inches(2.2), Inches(1.5), Inches(0.02), c)
    mul(sl, x+Inches(0.15), Inches(2.4), Inches(3.5), Inches(3.0), items, 12, TXT, Pt(4))

tb(sl, MG, Inches(6.2), CW, Inches(0.2),
   "设计模式：责任链 — TopologySort 编排，Flux.concat 串联，supportsProfile 按需激活", 10, MUT)

# ════════════════════════  6  ReAct 状态机  ════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); fill(sl)
page_title(sl, "06", "ReAct 引擎：3 状态流式状态机",
    "DefaultReActEngine · int[] 数组 · handle() + concatWith(defer)")

# 状态图
cx = SW / 2
sw = Inches(3.2); sg = Inches(0.45)
sx = cx - (sw*3 + sg*2) / 2

states = [("0  BUFFERING", C[0], "缓存非文本\n等待 content/tool_calls"),
          ("1  RELAYING", C[1], "逐 token 透传\nSSE message 事件"),
          ("2  TOOLS\nDETECTED", C[2], "合并剩余 chunk\n提取 tool_calls")]
for i, (name, c, desc) in enumerate(states):
    x = sx + Inches(i * (3.2 + 0.45))
    rnd(sl, x, Inches(1.15), Inches(3.2), Inches(1.5), c)
    tb(sl, x+Inches(0.1), Inches(1.2), Inches(3.0), Inches(0.45), name, 20, W, True, al=PP_ALIGN.CENTER)
    tb(sl, x+Inches(0.1), Inches(1.7), Inches(3.0), Inches(0.55), desc, 13, W, al=PP_ALIGN.CENTER)
    if i < len(states)-1:
        cond = "content 出现" if i == 0 else "tool_calls 出现"
        ax = x + Inches(3.22)
        rnd(sl, ax, Inches(1.65), Inches(0.43), Inches(0.3), MUT, 0.5)
        tb(sl, ax+Inches(0.01), Inches(1.67), Inches(0.41), Inches(0.25), cond, 8, W, True, al=PP_ALIGN.CENTER)

rnd(sl, sx+Inches(0.3), Inches(2.8), Inches(2.5), Inches(0.3), MUT, 0.04)
tb(sl, sx+Inches(0.35), Inches(2.83), Inches(2.4), Inches(0.25), "← 无 tool_calls 回到 BUFFERING", 9, W, True, al=PP_ALIGN.CENTER)

# 简化说明
rnd(sl, MG, Inches(3.2), Inches(12.5), Inches(3.5), W)
rect(sl, MG, Inches(3.2), Inches(12.5), Inches(0.04), PRI)
tb(sl, MG+Inches(0.12), Inches(3.28), Inches(5), Inches(0.22), "工作原理", 14, DK, True)
mul(sl, MG+Inches(0.12), Inches(3.6), Inches(12.2), Inches(2.8), [
    "▸ int[] state = {0} 记录三种状态：0=buffer / 1=relay / 2=tools",
    "▸ model.stream().handle() 逐 chunk 判别 → sink.next() 发射事件",
    "▸ 流结束后 concatWith(defer) 收尾：工具调用则进入多轮循环",
    "▸ content + tool_calls 可共存于同一 chunk，兼容混合模式",
    "▸ concatMap 保证工具串行执行，maxToolRounds=100 防无限循环",
    "▸ SSE 事件：THINKING / MESSAGE / TOOL_CALL / TOOL_APPROVAL / ERROR",
], 11, TXT, Pt(2))

tb(sl, MG, Inches(7.2), CW, Inches(0.2),
   "设计模式：状态机 — int 数组 + handle/concatWith 实现", 10, MUT)

# ════════════════════════  7  ReAct 多轮+审批  ════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); fill(sl)
page_title(sl, "06", "ReAct 引擎：多轮循环 & 审批机制",
    "handle + concatWith · concatMap 串行 · ApprovalStore")

# 顶部简化流程图 - 3 步 (整体下移 1.0 居中)
rnd(sl, MG, Inches(2.0), CW, Inches(0.7), W)
rect(sl, MG, Inches(2.0), CW, Inches(0.04), PRI)
tb(sl, MG+Inches(0.1), Inches(2.05), Inches(3), Inches(0.2), "多轮推理循环", 14, DK, True)
loop = [("执行工具", C[2]), ("结果回注", C[1]), ("继续推理", C[0])]
for i, (lbl, c) in enumerate(loop):
    x = MG + Inches(0.8 + i * 3.8)
    rect(sl, x, Inches(2.35), Inches(3.2), Inches(0.38), c, 0.04)
    tb(sl, x+Inches(0.04), Inches(2.37), Inches(3.12), Inches(0.34), lbl, 12, W, True, al=PP_ALIGN.CENTER)
    if i < len(loop)-1:
        tb(sl, x+Inches(3.25), Inches(2.33), Inches(0.5), Inches(0.3), "→", 22, MUT, True, al=PP_ALIGN.CENTER)
tb(sl, Inches(11.5), Inches(2.35), Inches(1.3), Inches(0.18), "↻ 100 轮上限", 10, MUT, True)

# 并排两张卡片 - 精简
card(sl, MG, Inches(3.0), Inches(6.0), Inches(2.8), PRI, "核心流程", [
    "▸ model.stream().handle() 逐 chunk 检测",
    "▸ 发现 tool_calls → concatMap 串行执行",
    "▸ 回注结果 → continueRounds 下一轮",
    "▸ 无工具 → 直接透传 SSE 事件",
], 13, 12)

card(sl, Inches(6.5), Inches(3.0), Inches(6.2), Inches(2.8), GRN, "审批保护", [
    "▸ @Tool(readonly=false) 标记需审批工具",
    "▸ CompletableFuture<Boolean> 异步等待",
    "▸ SSE → 前端 60s → 批准/拒绝/超时",
], 13, 12)

# ════════════════════════  8  反思拓扑——原语  ════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); fill(sl)
page_title(sl, "07", "反思拓扑：6 原语 & DAG 执行引擎",
    "TopologyExecutor · PrimitiveType · 原语可复用 · 拓扑可预设")

dag = [
    ("ACTOR", "入口：触发 LLM 推理\n建 ReflectionContext", C[0], MG, Inches(1.1)),
    ("EVALUATOR", "评分：4 维评估\n驱动路由决策", C[1], Inches(4.5), Inches(1.1)),
    ("REFLECTOR", "反思：对比期望\n生成改进建议", C[2], Inches(8.5), Inches(1.1)),
    ("SYNTHESIZER", "汇聚：合并多路\n结果 → 最终输出", C[3], Inches(8.5), Inches(2.7)),
    ("MEMORY", "记忆：跨轮存储\n推理时检索经验", C[4], Inches(4.5), Inches(2.7)),
    ("ROUTER", "路由：≥阈值走A\n<阈值走B", C[5], MG, Inches(2.7)),
]
for lbl, sub, c, x, y in dag:
    rnd(sl, x, y, Inches(3.8), Inches(1.4), c)
    tb(sl, x+Inches(0.1), y+Inches(0.12), Inches(3.6), Inches(0.35), lbl, 17, W, True, al=PP_ALIGN.CENTER)
    tb(sl, x+Inches(0.1), y+Inches(0.55), Inches(3.6), Inches(0.6), sub, 11, W, al=PP_ALIGN.CENTER)

arr(sl, Inches(3.85), Inches(1.55), Inches(0.6), Inches(0.3), MUT)
arr(sl, Inches(8.35), Inches(1.55), Inches(0.6), Inches(0.3), MUT)
arr(sl, Inches(3.85), Inches(3.15), Inches(0.6), Inches(0.3), MUT)
arr(sl, Inches(8.35), Inches(3.15), Inches(0.6), Inches(0.3), MUT)
tb(sl, Inches(3.2), Inches(2.8), Inches(1.2), Inches(0.18), "← 轮次循环", 9, MUT, True)
tb(sl, Inches(8.0), Inches(3.45), Inches(1.2), Inches(0.18), "→ 输出", 9, MUT, True)

# 扩展机制
rnd(sl, MG, Inches(4.5), CW, Inches(2.5), W)
rect(sl, MG, Inches(4.5), CW, Inches(0.04), PRI)
tb(sl, MG+Inches(0.12), Inches(4.58), Inches(5), Inches(0.22), "2 扩展机制", 14, DK, True)
mul(sl, MG+Inches(0.12), Inches(4.9), CW-Inches(0.25), Inches(1.8), [
    "▸ FORK 边类型 (EdgeType): 虚拟线程并行执行多个子图，CompletableFuture.allOf 汇合",
    "▸ COMPOSITE 原语: 嵌入子拓扑 (subTopology)，实现递归嵌套的复杂推理图",
], 11, TXT, Pt(4))

# 底部执行流程
rnd(sl, MG, Inches(6.2), CW, Inches(1.0), W)
rect(sl, MG, Inches(6.2), CW, Inches(0.04), BRO)
tb(sl, MG+Inches(0.12), Inches(6.3), Inches(10), Inches(0.18), "执行核心：TopologyExecutor.execute() — 从 entryNode 出发，沿边遍历 DAG", 11, DK, True)
tb(sl, MG+Inches(0.12), Inches(6.6), CW-Inches(0.25), Inches(0.4),
   "Actor → Evaluator → Router → Reflector → ... → Synthesizer → ExecutionResult.getFinalOutput()", 10, MUT)

tb(sl, MG, Inches(7.2), CW, Inches(0.2), "设计模式：DAG + 策略 — 节点可复用，边条件可配置，拓扑可预设", 10, MUT)

# ════════════════════════  9  反思拓扑——评估器 & 预设  ════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); fill(sl)
page_title(sl, "07", "反思拓扑：评估器 & 拓扑预设",
    "6 评估器 · CompositeEvaluator · 6 预设 · TopologyEvent 管线")

# 评估器表
rnd(sl, MG, Inches(1.0), Inches(6.0), Inches(3.0), W)
rect(sl, MG, Inches(1.0), Inches(6.0), Inches(0.04), GRN)
tb(sl, MG+Inches(0.1), Inches(1.08), Inches(3), Inches(0.2), "6 种评估器", 13, DK, True)
evals = [("LLMJudge", "LLM 语义级评估", PRI),
         ("Heuristic", "规则：长度/关键词", TEA),
         ("Consistency", "历史一致性校验", AMB),
         ("Importance", "信息重要性分级", GRN),
         ("ToolVerifier", "工具结果验证", BRO),
         ("Composite", "加权聚合多评估器", C[0])]
for i, (name, desc, c) in enumerate(evals):
    y = Inches(1.45 + i * 0.35)
    tag(sl, MG+Inches(0.12), y, Inches(1.2), Inches(0.25), c, name)
    tb(sl, MG+Inches(1.45), y+Inches(0.02), Inches(4.3), Inches(0.2), desc, 10, TXT)

# 4维评分
card(sl, Inches(7.0), Inches(1.0), Inches(2.8), Inches(1.6), PRI, "4 维评分", [
    "relevance · correctness",
    "completeness · clarity",
    "综合 ≥ 0.7 通过 / ROUTER 选路径",
], 12, 10)

# TopologyEvent 管线
card(sl, Inches(10.2), Inches(1.0), Inches(2.5), Inches(1.6), AMB, "事件管线", [
    "Sinks.Many 替代 Flux.create",
    "消除背压死锁",
    "reflect_step SSE → 前端",
    "MESSAGE + DONE 结束",
], 12, 10)

# 拓扑预设
rnd(sl, MG, Inches(2.9), CW, Inches(4.0), W)
rect(sl, MG, Inches(2.9), CW, Inches(0.04), TEA)
tb(sl, MG+Inches(0.12), Inches(2.98), Inches(3), Inches(0.2), "6 种拓扑预设", 14, DK, True)

presets = [
    ("passthrough", "ACTOR → EVAL → SYNTH (单轮直通)"),
    ("reflexion", "ACTOR → EVAL → ROUTER → REFLECTOR → ACTOR (反思循环)"),
    ("self-refine", "ACTOR → EVAL → ROUTER → ACTOR (自修正，无独立 Reflector)"),
    ("critic", "ACTOR → ToolVerifier → ROUTER → REFLECTOR → ACTOR"),
    ("multi-evaluator", "ACTOR → heuristic → llmJudge → ROUTER → REFLECTOR"),
    ("memory-augmented", "ACTOR → EVAL → ROUTER → [RETRY] MEMORY → REFLECTOR → ACTOR"),
]
for i, (name, desc) in enumerate(presets):
    y = Inches(3.35 + i * 0.5)
    tag(sl, MG+Inches(0.15), y, Inches(1.4), Inches(0.3), C[i % len(C)], name)
    tb(sl, MG+Inches(1.7), y+Inches(0.02), Inches(10.5), Inches(0.25), desc, 11, TXT)

tb(sl, MG, Inches(7.2), CW, Inches(0.2),
   "设计模式：策略 + 组合 — Evaluator 可互换，Composite 加权聚合多维度", 10, MUT)

# ════════════════════════  10  工具系统  ════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); fill(sl)
page_title(sl, "08", "工具系统：@Tool → 适配器 → 沙箱",
    "AnnotatedToolAdapter · 3 级沙箱 · 审批保护")

# 三阶段
phases = [("Phase1 扫描", C[0], "@Tool 注解收集"),
          ("Phase2 适配", C[1], "AnnotatedToolAdapter"),
          ("Phase3 注册", C[2], "ToolRegistry")]
for i, (lbl, c, desc) in enumerate(phases):
    x = MG + Inches(i * 4.0)
    rnd(sl, x, Inches(1.0), Inches(3.6), Inches(0.7), c)
    tb(sl, x+Inches(0.1), Inches(1.02), Inches(3.4), Inches(0.3), lbl, 13, W, True, al=PP_ALIGN.CENTER)
    tb(sl, x+Inches(0.1), Inches(1.32), Inches(3.4), Inches(0.2), desc, 11, W, al=PP_ALIGN.CENTER)
    if i < len(phases)-1: arr(sl, x+Inches(3.64), Inches(1.15), Inches(0.34), Inches(0.25), MUT)

# 工具表 (左)
th = ["工具", "级别", "超时"]
td = [("calculator", "DIRECT", "30s"), ("web_search", "DIRECT", "30s"),
      ("command", "PROCESS", "120s"), ("execute_script", "PROCESS", "30s")]
t = sl.shapes.add_table(5, 3, MG+Inches(0.1), Inches(2.0), Inches(4.3), Inches(1.5)).table
for ci, h in enumerate(th):
    c = t.cell(0, ci); c.text = h
    for p in c.text_frame.paragraphs:
        p.font.size = Pt(10); p.font.bold = True; p.font.color.rgb = W
    c.fill.solid(); c.fill.fore_color.rgb = PRI
for ri, row in enumerate(td):
    for ci, val in enumerate(row):
        c = t.cell(ri+1, ci); c.text = val
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(10); p.font.color.rgb = TXT
        c.fill.solid(); c.fill.fore_color.rgb = W if ri % 2 == 0 else BG

# 3 级沙箱 (右)
rnd(sl, Inches(7.0), Inches(2.0), Inches(5.7), Inches(1.5), W)
rect(sl, Inches(7.0), Inches(2.0), Inches(5.7), Inches(0.04), AMB)
tb(sl, Inches(7.15), Inches(2.08), Inches(3), Inches(0.2), "3 级沙箱隔离", 13, DK, True)
lvls = [("DIRECT", C[3], "当前线程直接执行"),
        ("SANDBOX", C[1], "守护线程 + 临时目录"),
        ("PROCESS", C[0], "OS 子进程 sh -c")]
for i, (lvl, c, ld) in enumerate(lvls):
    y = Inches(2.45 + i * 0.35)
    tag(sl, Inches(7.2), y, Inches(0.7), Inches(0.25), c, lvl)
    tb(sl, Inches(8.1), y+Inches(0.02), Inches(4.4), Inches(0.15), ld, 11, TXT)

# 底部全宽卡片 - 适配 + 审批 + 执行
rnd(sl, MG, Inches(3.8), CW, Inches(2.5), W)
rect(sl, MG, Inches(3.8), CW, Inches(0.04), GRN)
tb(sl, MG+Inches(0.12), Inches(3.88), Inches(6), Inches(0.2), "适配 + 审批 + 执行", 14, DK, True)
mul(sl, MG+Inches(0.12), Inches(4.2), CW-Inches(0.25), Inches(1.8), [
    "▸ AnnotatedToolAdapter: @Tool POJO → Tool 接口 (Adapter 模式)",
    "▸ 审批保护: @Tool(readonly=false) → CompletableFuture 异步等待前端确认",
    "▸ CommandExecutor: ProcessBuilder sh -c 子进程, 支持超时和截断",
    "▸ SSE 事件: tool_call 调用 / tool_approval 审批",
], 11, TXT, Pt(4))

tb(sl, MG, Inches(7.2), CW, Inches(0.2),
   "设计模式：适配器 (Adapter) + 外观 (Facade)", 10, MUT)

# ════════════════════════  11  未来演进  ════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); fill(sl)
page_title(sl, "09", "未来演进方向",
    "记忆架构 · 多 Agent 调度 · RAG 检索增强")

# 三个方向并排卡片
dirs = [
    ("记忆架构", BLU, "从无状态推理 → 有状态学习", [
        "4 层存储：工作 / episodic / 语义 / 程序",
        "MemoryHookRegistry 注入管线各阶段",
        "跨会话持久化 (SQLite/向量库)",
        "现有 memory-augmented 拓扑预设",
        "↓",
        "让 Agent 记住经验，避免重复错误",
    ]),
    ("多 Agent 调度", PNK, "从单 Agent → 多 Agent 协作", [
        "Supervisor + Worker 架构",
        "任务分解 → 派发 → 结果汇聚",
        "每个子 Agent 独立 Pipeline",
        "前端已支持子代理事件展示",
        "↓",
        "复杂任务拆解为多 Agent 协作",
    ]),
    ("RAG 检索增强", GRN, "从模型知识 → 私域知识增强", [
        "Embedding + 向量检索",
        "Pipeline 新增 ContextAugmentStage",
        "检索结果注入 SystemPrompt",
        "ReAct/Reflection 消费增强上下文",
        "↓",
        "私域问答/企业检索/代码理解",
    ]),
]
for i, (title, c, subtitle, items) in enumerate(dirs):
    x = MG + Inches(i * 4.1)
    rnd(sl, x, Inches(1.1), Inches(3.8), Inches(5.8), W)
    rect(sl, x, Inches(1.1), Inches(3.8), Inches(0.06), c)
    tb(sl, x+Inches(0.15), Inches(1.25), Inches(3.5), Inches(0.25), title, 18, c, True)
    tb(sl, x+Inches(0.15), Inches(1.55), Inches(3.5), Inches(0.2), subtitle, 11, MUT)
    rect(sl, x+Inches(0.15), Inches(1.8), Inches(1.5), Inches(0.02), c)
    mul(sl, x+Inches(0.15), Inches(2.0), Inches(3.5), Inches(4.5), items, 12, TXT, Pt(4))

# ════════════════════════  结尾 - 谢谢  ════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); fill(sl, DK)
rect(sl, 0, Inches(3.0), SW, Inches(0.03), PRI)
tb(sl, MG, Inches(2.5), CW, Inches(0.8), "谢谢！", 52, W, True, al=PP_ALIGN.CENTER)
tb(sl, MG, Inches(3.5), CW, Inches(0.4), "LyClaw — 通用 AI Agent 框架", 16, MUT, al=PP_ALIGN.CENTER)
rect(sl, 0, Inches(4.5), SW, Inches(0.03), PRI)

out = "/tmp/LyClaw-技术亮点串讲.pptx"
prs.save(out)
print(f"✅ PPT: {out}  ({len(prs.slides)} 页)")
